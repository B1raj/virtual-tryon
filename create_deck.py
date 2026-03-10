from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Brand palette ──────────────────────────────────────────────
NAVY       = RGBColor(0x0D, 0x1B, 0x3E)
ACCENT     = RGBColor(0x00, 0xC2, 0xFF)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xE8, 0xEC, 0xF0)
MID_GRAY   = RGBColor(0x8A, 0x99, 0xAA)
GOLD       = RGBColor(0xFF, 0xC1, 0x07)
DARK_CARD  = RGBColor(0x0A, 0x14, 0x30)
GREEN      = RGBColor(0x2E, 0xCC, 0x71)
RED_LIGHT  = RGBColor(0xFF, 0x6B, 0x6B)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]

# ══════════════════════════════════════════════════════════════
# HELPERS
# ══════════════════════════════════════════════════════════════

def rect(sl, l, t, w, h, fill):
    sh = sl.shapes.add_shape(1, l, t, w, h)
    sh.line.fill.background()
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    return sh

def txt(sl, text, l, t, w, h, size=16, bold=False,
        color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    tb = sl.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    r  = p.add_run()
    r.text           = text
    r.font.size      = Pt(size)
    r.font.bold      = bold
    r.font.italic    = italic
    r.font.color.rgb = color
    return tb

def bullets(sl, items, l, t, w, h, size=15,
            color=WHITE, bc=None, sub_items=None):
    bc = bc or ACCENT
    tb = sl.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        br = p.add_run(); br.text = "▸  "
        br.font.size = Pt(size); br.font.bold = True
        br.font.color.rgb = bc
        cr = p.add_run(); cr.text = item
        cr.font.size = Pt(size)
        cr.font.color.rgb = color

def divider(sl, t, color=ACCENT):
    line = sl.shapes.add_shape(1, Inches(0.5), t, Inches(12.33), Pt(2))
    line.fill.solid(); line.fill.fore_color.rgb = color
    line.line.fill.background()

def slide_chrome(sl, num_str, title_str):
    """Standard slide background + header."""
    rect(sl, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    rect(sl, 0, 0, Inches(0.18), SLIDE_H, ACCENT)
    rect(sl, Inches(0.18), 0, SLIDE_W, Inches(1.1), DARK_CARD)
    txt(sl, num_str, Inches(0.3), Inches(0.1), Inches(0.9), Inches(0.9),
        size=44, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT)
    txt(sl, title_str, Inches(1.3), Inches(0.25), Inches(10), Inches(0.65),
        size=30, bold=True, color=WHITE)
    divider(sl, Inches(1.12))

def stat_card(sl, val, lbl, l, t, w=Inches(2.8)):
    rect(sl, l, t, w, Inches(1.05), DARK_CARD)
    rect(sl, l, t, w, Inches(0.05), ACCENT)
    txt(sl, val, l, t + Inches(0.1), w, Inches(0.55),
        size=26, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    txt(sl, lbl, l, t + Inches(0.65), w, Inches(0.35),
        size=11, color=MID_GRAY, align=PP_ALIGN.CENTER)

def label(sl, text, l, t, w, h, color=ACCENT):
    txt(sl, text, l, t, w, h, size=11, bold=True, color=color)

# ══════════════════════════════════════════════════════════════
# SLIDE 1 — Title / Hero
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
rect(s, 0, 0, Inches(0.18), SLIDE_H, ACCENT)
rect(s, Inches(0.18), 0, SLIDE_W, Inches(0.08), ACCENT)

# Course chip
rect(s, Inches(0.5), Inches(0.3), Inches(3.4), Inches(0.45), ACCENT)
txt(s, "MSDS · COMPUTERVISION 462  |  Final Project",
    Inches(0.55), Inches(0.3), Inches(3.3), Inches(0.45),
    size=9, bold=True, color=NAVY)

txt(s, "Enhanced Virtual Try-On\nin the Metaverse",
    Inches(0.5), Inches(1.1), Inches(8.5), Inches(2.2),
    size=48, bold=True)
txt(s, "Leveraging AI & Pretrained Models for High-Resolution\nClothing Try-On Experiences",
    Inches(0.5), Inches(3.1), Inches(8.5), Inches(1.2),
    size=20, color=ACCENT)
rect(s, Inches(0.5), Inches(4.45), Inches(5), Inches(0.04), MID_GRAY)
txt(s, "Team: Joyati  ·  Biraj Mishra  ·  Murughanandam S.",
    Inches(0.5), Inches(4.6), Inches(7), Inches(0.4),
    size=14, color=LIGHT_GRAY)
txt(s, "Spring 2026",
    Inches(0.5), Inches(5.0), Inches(3), Inches(0.4),
    size=12, color=MID_GRAY)
txt(s, "6G", Inches(9.8), Inches(1.5), Inches(3), Inches(3),
    size=160, bold=True, color=RGBColor(0x1A, 0x2D, 0x55),
    align=PP_ALIGN.CENTER)
txt(s, "⬡", Inches(10.3), Inches(3.6), Inches(2), Inches(1.5),
    size=80, color=RGBColor(0x00, 0x4A, 0x6E), align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# SLIDE 2 — Executive Summary
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_chrome(s, "01", "Executive Summary")

# Four pillars
pillars = [
    ("🛍️ BUSINESS PROBLEM",  GOLD,
     "Online shoppers can't try clothes before buying — leading to 30–40% return rates. "
     "Virtual try-on eliminates this friction, driving conversion and reducing logistics costs."),
    ("🤖 OUR SOLUTION",       ACCENT,
     "An AI-powered pipeline that digitally dresses a shopper in any garment using their photo. "
     "Built on HR-VITON — a state-of-the-art pretrained model from ECCV 2022."),
    ("📈 KEY RESULT",         GREEN,
     "SSIM score of 0.85 on paired evaluation — matching published research benchmarks. "
     "7× improvement over our trained-from-scratch baseline (SSIM 0.12)."),
    ("🚀 BUSINESS VALUE",     RGBColor(0xFF, 0x8C, 0x00),
     "Reduces product returns · Increases purchase confidence · Scales to any catalog size "
     "· Ready for integration into e-commerce and metaverse platforms."),
]

xs = [Inches(0.45), Inches(3.65), Inches(6.85), Inches(10.05)]
for i, ((heading, hcolor, body_text), x) in enumerate(zip(pillars, xs)):
    rect(s, x, Inches(1.3), Inches(3.0), Inches(5.8), DARK_CARD)
    rect(s, x, Inches(1.3), Inches(3.0), Inches(0.06), hcolor)
    txt(s, heading, x + Inches(0.15), Inches(1.45),
        Inches(2.7), Inches(0.5), size=10, bold=True, color=hcolor)
    txt(s, body_text, x + Inches(0.15), Inches(2.05),
        Inches(2.7), Inches(4.8), size=13, color=LIGHT_GRAY)

# ══════════════════════════════════════════════════════════════
# SLIDE 3 — Team Introduction
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_chrome(s, "02", "Meet the Team")

team = [
    {"name": "Joyati",           "role": "ML Engineer",
     "focus": "Model Testing & Evaluation",
     "icon": "🔬",
     "details": [
         "Model testing & monitoring strategy",
         "Evaluation pipelines & QA metrics",
         "Try-on synthesis quality assessment",
     ]},
    {"name": "Biraj Mishra",     "role": "Software Engineer",
     "focus": "Systems & Integration",
     "icon": "⚙️",
     "details": [
         "System architecture & API design",
         "CV model integration into pipelines",
         "End-to-end notebook engineering",
     ]},
    {"name": "Murughanandam S.", "role": "Senior Director",
     "focus": "Research & Strategy",
     "icon": "🎯",
     "details": [
         "Research direction & benchmarking",
         "Connects academia to real deployment",
         "Stakeholder communication & roadmap",
     ]},
]

for i, (m, lx) in enumerate(zip(team, [Inches(0.5), Inches(4.6), Inches(8.7)])):
    rect(s, lx, Inches(1.3), Inches(3.9), Inches(5.7), DARK_CARD)
    rect(s, lx, Inches(1.3), Inches(3.9), Inches(0.07), ACCENT)
    rect(s, lx + Inches(1.45), Inches(1.5), Inches(1.0), Inches(1.0),
         RGBColor(0x00, 0x4A, 0x6E))
    txt(s, m["icon"], lx + Inches(1.45), Inches(1.5), Inches(1.0), Inches(1.0),
        size=32, align=PP_ALIGN.CENTER)
    txt(s, m["name"], lx + Inches(0.1), Inches(2.65),
        Inches(3.7), Inches(0.45), size=16, bold=True, align=PP_ALIGN.CENTER)
    rect(s, lx + Inches(0.65), Inches(3.15), Inches(2.6), Inches(0.32), ACCENT)
    txt(s, m["role"], lx + Inches(0.65), Inches(3.15),
        Inches(2.6), Inches(0.32), size=10, bold=True, color=NAVY,
        align=PP_ALIGN.CENTER)
    txt(s, m["focus"].upper(), lx + Inches(0.15), Inches(3.58),
        Inches(3.6), Inches(0.28), size=9, bold=True, color=MID_GRAY,
        align=PP_ALIGN.CENTER)
    rect(s, lx + Inches(0.3), Inches(3.92), Inches(3.3), Inches(0.025), MID_GRAY)
    bullets(s, m["details"], lx + Inches(0.15), Inches(4.0),
            Inches(3.6), Inches(2.5), size=12, color=LIGHT_GRAY)

# ══════════════════════════════════════════════════════════════
# SLIDE 4 — Problem Statement
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_chrome(s, "03", "Problem Statement")

txt(s, "CONTEXT", Inches(0.5), Inches(1.25), Inches(5.8), Inches(0.3),
    size=11, bold=True, color=ACCENT)
txt(s, "Online fashion retail demands immersive, realistic try-on experiences — "
       "yet current virtual fitting rooms fall short on resolution, realism, and scale. "
       "Every failed try-on experience translates directly to lost revenue.",
    Inches(0.5), Inches(1.6), Inches(5.8), Inches(1.0),
    size=14, color=LIGHT_GRAY)

challenges = [
    "30–40% return rates cost retailers billions annually in reverse logistics",
    "Low-resolution outputs (256×192) create visible artifacts — customers notice",
    "Garment misalignment causes unnatural, unconvincing try-on results",
    "No scalable real-time solution exists for Metaverse & next-gen retail",
]
bullets(s, challenges, Inches(0.5), Inches(2.65), Inches(5.8), Inches(3.0), size=14)

# Right panel — focus areas
rect(s, Inches(6.8), Inches(1.2), Inches(6.2), Inches(5.9), DARK_CARD)
txt(s, "WHAT WE SET OUT TO SOLVE", Inches(7.0), Inches(1.38),
    Inches(5.8), Inches(0.35), size=11, bold=True, color=GOLD)

focus = [
    ("Photo-Realistic Try-On",
     "Generate convincing 768×1024 images of a person wearing any chosen garment."),
    ("Accurate Cloth Warping",
     "Precisely align the garment to the body's shape, pose, and proportions."),
    ("Scalable AI Pipeline",
     "Build a modular system that works across any catalog — no re-training needed."),
    ("Measurable Quality",
     "Use rigorous metrics (SSIM, PSNR) so results can be objectively benchmarked."),
]
y = Inches(1.9)
for title, desc in focus:
    rect(s, Inches(6.95), y, Inches(0.05), Inches(0.65), ACCENT)
    txt(s, title, Inches(7.15), y, Inches(5.5), Inches(0.35),
        size=14, bold=True, color=ACCENT)
    txt(s, desc, Inches(7.15), y + Inches(0.32), Inches(5.5), Inches(0.5),
        size=13, color=LIGHT_GRAY)
    y += Inches(1.2)

# ══════════════════════════════════════════════════════════════
# SLIDE 5 — Data Profile
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_chrome(s, "04", "Data Profile — VITON-HD Dataset")

rect(s, Inches(0.45), Inches(1.25), Inches(8.0), Inches(0.52), DARK_CARD)
txt(s, "VITON-HD  ·  High-Resolution Zalando Dataset  (publicly available on Kaggle)",
    Inches(0.55), Inches(1.27), Inches(7.8), Inches(0.48),
    size=14, bold=True, color=ACCENT)

txt(s, "A paired person–garment benchmark enabling supervised learning. "
       "Each sample provides the full context needed to teach an AI model "
       "how to place and blend a garment onto a person realistically.",
    Inches(0.45), Inches(1.85), Inches(7.8), Inches(0.85),
    size=13, color=LIGHT_GRAY)

# Stats
for j, (val, lbl) in enumerate([
    ("11,647", "Training Pairs"), ("2,032", "Test Pairs"),
    ("768×1024", "Image Resolution"), ("9", "Data Modalities"),
]):
    stat_card(s, val, lbl, Inches(0.45) + j * Inches(3.05), Inches(2.9), Inches(2.8))

# Modalities
txt(s, "WHAT'S IN EACH SAMPLE", Inches(0.45), Inches(4.15),
    Inches(7.8), Inches(0.3), size=11, bold=True, color=GOLD)

mods = [
    ("Person Photo", "High-res image of the model"),
    ("Target Garment", "Flat-lay photo of the clothing item"),
    ("Cloth Mask", "Precise garment outline for compositing"),
    ("Body Parse Map", "18-region body segmentation (hair, arms, torso…)"),
    ("Agnostic Image", "Person photo with clothing region blanked out"),
    ("Body Keypoints", "25-point skeleton from OpenPose for pose alignment"),
    ("DensePose Map", "3D surface UV coordinates mapped onto the body"),
    ("Parse-Agnostic", "Segmentation map with clothing labels removed"),
    ("OpenPose Image", "Rendered skeleton overlay for visual debugging"),
]

col_w = Inches(3.85)
for i, (name, desc) in enumerate(mods):
    col = i % 2
    row = i // 2
    lx = Inches(0.45) + col * Inches(4.05)
    ty = Inches(4.55) + row * Inches(0.52)
    if ty > Inches(7.0): break
    txt(s, f"▸  {name}: ", lx, ty, Inches(1.55), Inches(0.45),
        size=12, bold=True, color=ACCENT)
    txt(s, desc, lx + Inches(1.5), ty, col_w - Inches(1.5), Inches(0.45),
        size=12, color=LIGHT_GRAY)

# Right panel
rect(s, Inches(8.7), Inches(1.2), Inches(4.3), Inches(6.0), DARK_CARD)
txt(s, "WHY PAIRED DATA MATTERS", Inches(8.85), Inches(1.38),
    Inches(4.0), Inches(0.35), size=11, bold=True, color=ACCENT)
txt(s, "Paired training means the AI can directly compare its output "
       "against the real photo — like a teacher grading with an answer key. "
       "This enables measurable, objective improvement.",
    Inches(8.85), Inches(1.82), Inches(4.0), Inches(1.2),
    size=12, color=LIGHT_GRAY)
txt(s, "TRAIN / TEST SPLIT", Inches(8.85), Inches(3.15),
    Inches(4.0), Inches(0.3), size=11, bold=True, color=GOLD)
for lbl, pct, fill in [("Train  85%", 85, ACCENT), ("Test    15%", 15, MID_GRAY)]:
    pass  # visual bar approximation
bar_y = Inches(3.5)
for lbl, pct_w, fc in [
    ("Train  85%  (11,647 pairs)", Inches(3.4), ACCENT),
    ("Test   15%  ( 2,032 pairs)", Inches(0.6), MID_GRAY),
]:
    rect(s, Inches(8.85), bar_y, pct_w, Inches(0.4), fc)
    txt(s, lbl, Inches(8.85), bar_y + Inches(0.44), Inches(4.0), Inches(0.3),
        size=10, color=LIGHT_GRAY)
    bar_y += Inches(0.9)

txt(s, "PREPROCESSING STEPS", Inches(8.85), Inches(5.45),
    Inches(4.0), Inches(0.3), size=11, bold=True, color=ACCENT)
bullets(s, [
    "Resize to model input (256×192 or 768×1024)",
    "Normalise pixel values to [-1, 1]",
    "Horizontal flip + colour jitter for variety",
], Inches(8.85), Inches(5.8), Inches(4.0), Inches(1.4), size=11)

# ══════════════════════════════════════════════════════════════
# SLIDE 6 — Exploratory Data Analysis
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_chrome(s, "05", "Exploratory Data Analysis — Key Findings")

txt(s, "Before building models, we studied the dataset to understand its patterns "
       "and inform our design choices.",
    Inches(0.5), Inches(1.2), Inches(12.5), Inches(0.45),
    size=13, color=MID_GRAY, italic=True)

findings = [
    {
        "title": "🖼️  Image Statistics",
        "color": ACCENT,
        "points": [
            "Person images: mean brightness ~128, garments slightly brighter (~140)",
            "RGB channels consistent across dataset — no significant colour bias",
            "Resolution always 768×1024 — no resizing artifacts at native size",
        ],
        "insight": "Insight: Normalising to [-1,1] is safe; no outlier images.",
    },
    {
        "title": "🎨  Body Parse Map Analysis",
        "color": GOLD,
        "points": [
            "Background: ~50% of pixels (expected — most of image is white)",
            "Upper-clothes (label 5): ~15% — the primary training target region",
            "Hair, face, arms together: ~20% — must be preserved in output",
        ],
        "insight": "Insight: Model must focus on label 5 without disturbing other regions.",
    },
    {
        "title": "📐  Cloth Mask Coverage",
        "color": GREEN,
        "points": [
            "Mean garment coverage: ~30% of image area",
            "Standard deviation: ±8% — garments vary in size",
            "Bounding-box aspect ratio: mean 0.72 (wider than tall — typical shirts)",
        ],
        "insight": "Insight: Warp model must handle varying garment sizes gracefully.",
    },
    {
        "title": "🦴  Body Keypoint Quality",
        "color": RGBColor(0xFF, 0x8C, 0x00),
        "points": [
            "OpenPose confidence > 0.7 for 85% of upper-body keypoints",
            "Neck & shoulder keypoints most reliable — ideal warp anchors",
            "Wrist keypoints occasionally missing (~15% of samples)",
        ],
        "insight": "Insight: Upper-body keypoints are reliable; wrist fallback needed.",
    },
]

xs = [Inches(0.45), Inches(3.65), Inches(6.85), Inches(10.05)]
for f, lx in zip(findings, xs):
    rect(s, lx, Inches(1.7), Inches(3.0), Inches(5.5), DARK_CARD)
    rect(s, lx, Inches(1.7), Inches(3.0), Inches(0.06), f["color"])
    txt(s, f["title"], lx + Inches(0.12), Inches(1.82),
        Inches(2.76), Inches(0.55), size=11, bold=True, color=f["color"])
    bullets(s, f["points"], lx + Inches(0.12), Inches(2.45),
            Inches(2.76), Inches(2.8), size=10.5, color=LIGHT_GRAY,
            bc=f["color"])
    rect(s, lx, Inches(5.95), Inches(3.0), Inches(0.03), f["color"])
    txt(s, f["insight"], lx + Inches(0.12), Inches(6.02),
        Inches(2.76), Inches(0.9), size=10, color=f["color"], italic=True)

# ══════════════════════════════════════════════════════════════
# SLIDE 7 — Technology Platform & Architecture
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_chrome(s, "06", "Technology Platform & Solution Architecture")

# Left: stack
rect(s, Inches(0.45), Inches(1.25), Inches(4.8), Inches(5.95), DARK_CARD)
txt(s, "TECHNOLOGY STACK", Inches(0.6), Inches(1.4),
    Inches(4.5), Inches(0.3), size=11, bold=True, color=ACCENT)

stack = [
    ("Core Framework",  "PyTorch 2.2 + torchvision 0.17"),
    ("Model",           "HR-VITON (ECCV 2022 official weights)"),
    ("Image Processing","OpenCV 4.12 · Pillow 10.3"),
    ("Scientific",      "NumPy · SciPy · scikit-image"),
    ("Visualisation",   "Matplotlib 3.9 · Seaborn"),
    ("Notebooks",       "Jupyter · nbformat 5.10"),
    ("Metrics",         "SSIM · PSNR · MSE (scikit-image)"),
    ("Cloud / Compute", "Google Colab GPU (T4) + AMP FP16"),
    ("Dataset Source",  "Kaggle — Zalando VITON-HD"),
    ("Code Repo",       "GitHub — version-controlled notebooks"),
]
y = Inches(1.85)
for cat, val in stack:
    txt(s, cat, Inches(0.6), y, Inches(1.7), Inches(0.35),
        size=10.5, bold=True, color=ACCENT)
    txt(s, val, Inches(2.3), y, Inches(2.8), Inches(0.35),
        size=10.5, color=LIGHT_GRAY)
    y += Inches(0.5)

# Centre: pipeline diagram
rect(s, Inches(5.55), Inches(1.25), Inches(4.8), Inches(5.95), DARK_CARD)
txt(s, "INFERENCE PIPELINE", Inches(5.7), Inches(1.4),
    Inches(4.5), Inches(0.3), size=11, bold=True, color=GOLD)

steps = [
    (ACCENT,              "INPUT",       "Person photo + target garment"),
    (RGBColor(0x00,0x8C,0xB8), "STAGE 1","ConditionGenerator — warp garment to body"),
    (RGBColor(0x00,0x6A,0x8A), "STAGE 2","SPADEGenerator — synthesise final image"),
    (GREEN,               "OUTPUT",      "Photo-realistic try-on result (768×1024)"),
    (MID_GRAY,            "EVALUATE",    "SSIM / PSNR / MSE vs. ground truth"),
]
dy = Inches(1.9)
for fc, stage, desc in steps:
    rect(s, Inches(5.7), dy, Inches(4.4), Inches(0.55), fc)
    txt(s, stage, Inches(5.75), dy + Inches(0.04),
        Inches(1.1), Inches(0.45), size=10, bold=True, color=NAVY)
    txt(s, desc, Inches(6.85), dy + Inches(0.04),
        Inches(3.2), Inches(0.45), size=11, color=NAVY, bold=True)
    if stage != "EVALUATE":
        txt(s, "▼", Inches(7.7), dy + Inches(0.56), Inches(0.5), Inches(0.3),
            size=10, color=MID_GRAY, align=PP_ALIGN.CENTER)
    dy += Inches(0.9)

# Right: deployment
rect(s, Inches(10.6), Inches(1.25), Inches(2.45), Inches(5.95), DARK_CARD)
txt(s, "DEPLOYMENT", Inches(10.72), Inches(1.4),
    Inches(2.2), Inches(0.3), size=11, bold=True, color=ACCENT)

deploy_items = [
    ("📦", "Modular Python\npackage"),
    ("☁️", "Cloud GPU\n(Colab / AWS)"),
    ("🔌", "REST API\nready"),
    ("🛍️", "E-commerce\nintegration"),
    ("⚡", "~2s per\ntry-on image"),
]
dy = Inches(1.85)
for icon, desc in deploy_items:
    rect(s, Inches(10.72), dy, Inches(2.1), Inches(0.72), RGBColor(0x12,0x20,0x48))
    txt(s, icon, Inches(10.75), dy, Inches(0.45), Inches(0.72),
        size=18, align=PP_ALIGN.CENTER)
    txt(s, desc, Inches(11.22), dy + Inches(0.03),
        Inches(1.55), Inches(0.65), size=10, color=LIGHT_GRAY)
    dy += Inches(0.85)

# ══════════════════════════════════════════════════════════════
# SLIDE 8 — Methodology
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_chrome(s, "07", "Methodology — How We Built the Solution")

txt(s, "Our approach evolved through three versions, each addressing failures of the last.",
    Inches(0.5), Inches(1.2), Inches(12.5), Inches(0.4),
    size=13, color=MID_GRAY, italic=True)

versions = [
    {
        "ver": "v2",
        "title": "CP-VTON  ·  Train from Scratch",
        "color": RED_LIGHT,
        "status": "❌  SSIM 0.12",
        "what": [
            "Built a classic two-stage pipeline from scratch",
            "Stage 1: CNN feature matching → TPS warp",
            "Stage 2: U-Net synthesis with VGG perceptual loss",
        ],
        "why_failed": [
            "Buggy TPS warp kernel → distorted garments",
            "Hardcoded CPU → 8 min/epoch, only 600 pairs",
            "Missing VGG normalisation → meaningless loss",
        ],
    },
    {
        "ver": "v3",
        "title": "Dense Flow  ·  Fixed & Upgraded",
        "color": GOLD,
        "status": "⚠️  SSIM 0.12",
        "what": [
            "Replaced buggy TPS with dense flow UNet",
            "Auto GPU + mixed precision (AMP FP16)",
            "5,000 training pairs; ALIAS decoder blocks",
        ],
        "why_failed": [
            "Architecture fixes weren't enough alone",
            "Training from scratch needs 100k+ pairs",
            "Colab compute budget = ~10 GPU-hours",
        ],
    },
    {
        "ver": "v4",
        "title": "HR-VITON  ·  Pretrained Model",
        "color": GREEN,
        "status": "✅  SSIM 0.85",
        "what": [
            "Load official ECCV 2022 pretrained weights",
            "Run inference only — no training needed",
            "Stage 1: ConditionGenerator warps cloth",
            "Stage 2: SPADEGenerator renders final image",
        ],
        "why_failed": None,
    },
]

xs = [Inches(0.45), Inches(4.65), Inches(8.85)]
ws = [Inches(4.0),  Inches(4.0),  Inches(4.15)]
for v, lx, cw in zip(versions, xs, ws):
    rect(s, lx, Inches(1.65), cw, Inches(5.55), DARK_CARD)
    rect(s, lx, Inches(1.65), cw, Inches(0.07), v["color"])
    txt(s, v["ver"], lx + Inches(0.15), Inches(1.78),
        Inches(0.6), Inches(0.42), size=22, bold=True, color=v["color"])
    txt(s, v["title"], lx + Inches(0.8), Inches(1.82),
        cw - Inches(0.95), Inches(0.38), size=12, bold=True, color=WHITE)
    rect(s, lx + cw - Inches(1.45), Inches(1.75),
         Inches(1.4), Inches(0.32), RGBColor(0x12,0x20,0x48))
    txt(s, v["status"], lx + cw - Inches(1.45), Inches(1.76),
        Inches(1.4), Inches(0.3), size=10, bold=True, color=v["color"],
        align=PP_ALIGN.CENTER)
    txt(s, "APPROACH", lx + Inches(0.15), Inches(2.3),
        cw - Inches(0.3), Inches(0.25), size=9, bold=True, color=ACCENT)
    bullets(s, v["what"], lx + Inches(0.15), Inches(2.6),
            cw - Inches(0.3), Inches(2.0), size=11, color=LIGHT_GRAY)
    if v["why_failed"]:
        txt(s, "WHY IT FELL SHORT", lx + Inches(0.15), Inches(4.5),
            cw - Inches(0.3), Inches(0.25), size=9, bold=True, color=RED_LIGHT)
        bullets(s, v["why_failed"], lx + Inches(0.15), Inches(4.8),
                cw - Inches(0.3), Inches(1.5), size=11,
                color=LIGHT_GRAY, bc=RED_LIGHT)
    else:
        txt(s, "KEY INSIGHT", lx + Inches(0.15), Inches(4.5),
            cw - Inches(0.3), Inches(0.25), size=9, bold=True, color=GREEN)
        txt(s, "Training from scratch requires massive data & compute. "
               "A pretrained model delivers research-grade quality instantly — "
               "the pragmatic choice for production.",
            lx + Inches(0.15), Inches(4.8), cw - Inches(0.3), Inches(1.4),
            size=11, color=GREEN, italic=True)

# Arrow between cards
for ax in [Inches(4.52), Inches(8.72)]:
    txt(s, "→", ax, Inches(4.1), Inches(0.25), Inches(0.4),
        size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# SLIDE 9 — Literature Review
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_chrome(s, "08", "Literature Review & Academic Foundation")

papers = [
    {
        "ref": "VITON  (Han et al., CVPR 2018)",
        "color": ACCENT,
        "desc": "The original image-based virtual try-on paper. "
                "Introduced the two-stage paradigm: geometric matching → synthesis. "
                "Worked at 256×192 — quality insufficient for commercial use.",
        "our_use": "Foundation architecture. We started here for v2.",
    },
    {
        "ref": "CP-VTON  (Wang et al., ECCV 2018)",
        "color": RGBColor(0x00, 0xA0, 0xD0),
        "desc": "Added Thin-Plate Spline (TPS) warping to the GMM stage for "
                "smoother cloth deformation. Introduced VGG perceptual loss. "
                "Still limited to 256×192 resolution.",
        "our_use": "Direct blueprint for our v2 implementation.",
    },
    {
        "ref": "VITON-HD  (Choi et al., CVPR 2021)",
        "color": GOLD,
        "desc": "First high-resolution (1024×768) try-on model. "
                "Introduced ALIAS (Adaptive Local Instance-Aware) normalisation "
                "to handle misalignment between warped cloth and body.",
        "our_use": "ALIAS concept adopted in our v3 TOM decoder.",
    },
    {
        "ref": "HR-VITON  (Lee et al., ECCV 2022)",
        "color": GREEN,
        "desc": "State of the art. Jointly predicts cloth warping and body "
                "segmentation in one network (ConditionGenerator). "
                "SPADEGenerator produces sharp 768×1024 outputs. "
                "Published SSIM 0.844 / PSNR 27.1 dB.",
        "our_use": "Our v4 — loaded pretrained weights directly.",
    },
]

xs = [Inches(0.45), Inches(3.65), Inches(6.85), Inches(10.05)]
for p, lx in zip(papers, xs):
    rect(s, lx, Inches(1.3), Inches(3.0), Inches(5.85), DARK_CARD)
    rect(s, lx, Inches(1.3), Inches(3.0), Inches(0.06), p["color"])
    txt(s, p["ref"], lx + Inches(0.12), Inches(1.42),
        Inches(2.76), Inches(0.55), size=10.5, bold=True, color=p["color"])
    txt(s, p["desc"], lx + Inches(0.12), Inches(2.05),
        Inches(2.76), Inches(3.0), size=11, color=LIGHT_GRAY)
    rect(s, lx, Inches(5.75), Inches(3.0), Inches(0.03), p["color"])
    txt(s, "OUR USE:", lx + Inches(0.12), Inches(5.83),
        Inches(0.85), Inches(0.3), size=9, bold=True, color=p["color"])
    txt(s, p["our_use"], lx + Inches(0.97), Inches(5.83),
        Inches(1.9), Inches(0.5), size=9.5, color=p["color"], italic=True)

txt(s, "Additional references: U-Net (Ronneberger, MICCAI 2015) · "
       "OpenPose (Cao et al., TPAMI 2019) · SPADE (Park et al., CVPR 2019) · "
       "Adam optimizer (Kingma & Ba, ICLR 2015) · SSIM (Wang et al., IEEE TIP 2004)",
    Inches(0.45), Inches(7.1), Inches(12.5), Inches(0.3),
    size=9, color=MID_GRAY, italic=True)

# ══════════════════════════════════════════════════════════════
# SLIDE 10 — Results & Evaluation
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_chrome(s, "09", "Results & Evaluation")

# Results table
headers = ["Method", "SSIM ↑", "PSNR ↑", "vs. Baseline"]
rows = [
    ("No-change baseline (agnostic image)", "0.45", "~18 dB", "—",         MID_GRAY,   False),
    ("Classical warp (bounding-box, no AI)","0.52", "~19 dB", "+16%",      LIGHT_GRAY, False),
    ("v2 — CP-VTON trained from scratch",   "0.12", "~14 dB", "−73%",      RED_LIGHT,  False),
    ("v3 — Dense Flow trained (Colab)",     "0.12", "~14 dB", "−73%",      RED_LIGHT,  False),
    ("v4 — HR-VITON pretrained (paired)",   "0.85", "~27 dB", "+89%",      GREEN,      True),
    ("v4 — HR-VITON pretrained (unpaired)", "0.70", "~22 dB", "+56%",      ACCENT,     False),
    ("Published HR-VITON (Lee et al. 2022)","0.844","27.1 dB","benchmark", GOLD,       False),
]

col_xs = [Inches(0.45), Inches(6.5), Inches(8.3), Inches(10.2)]
col_ws = [Inches(5.9),  Inches(1.65),Inches(1.75),Inches(1.95)]
th = Inches(1.28)

# Header row
rect(s, Inches(0.45), th, Inches(12.2), Inches(0.35), ACCENT)
for hdr, lx, lw in zip(headers, col_xs, col_ws):
    align = PP_ALIGN.LEFT if hdr == "Method" else PP_ALIGN.CENTER
    txt(s, hdr, lx + Inches(0.07), th + Inches(0.02),
        lw - Inches(0.1), Inches(0.3), size=11, bold=True, color=NAVY, align=align)

row_h = Inches(0.47)
for i, (method, ssim, psnr, vs, color, highlight) in enumerate(rows):
    ry = th + Inches(0.35) + i * row_h
    bg = RGBColor(0x08,0x1A,0x3A) if highlight else (
         DARK_CARD if i % 2 == 0 else RGBColor(0x0E,0x1E,0x42))
    rect(s, Inches(0.45), ry, Inches(12.2), row_h - Inches(0.02), bg)
    if highlight:
        rect(s, Inches(0.45), ry, Inches(0.05), row_h - Inches(0.02), GREEN)
    txt(s, method, col_xs[0] + Inches(0.12), ry + Inches(0.06),
        col_ws[0] - Inches(0.15), Inches(0.35), size=11, color=color, bold=highlight)
    for val, lx, lw in [(ssim, col_xs[1], col_ws[1]),
                         (psnr, col_xs[2], col_ws[2]),
                         (vs,   col_xs[3], col_ws[3])]:
        txt(s, val, lx, ry + Inches(0.06), lw, Inches(0.35),
            size=11, color=color, bold=highlight, align=PP_ALIGN.CENTER)

# Key takeaways
rect(s, Inches(0.45), Inches(5.65), Inches(12.2), Inches(1.5), DARK_CARD)
txt(s, "KEY TAKEAWAYS", Inches(0.6), Inches(5.75),
    Inches(4.0), Inches(0.28), size=10, bold=True, color=GOLD)
takeaways = [
    "Pretrained HR-VITON (v4) achieves SSIM 0.85 — matching the published research paper",
    "7× improvement over our trained-from-scratch attempts (SSIM 0.12)",
    "Unpaired SSIM 0.70 is expected — the output correctly shows a different garment, not a failure",
    "Training from scratch requires 100k+ samples & weeks of GPU time — beyond course scope",
]
bullets(s, takeaways, Inches(0.6), Inches(6.08), Inches(12.0), Inches(1.0),
        size=11, color=LIGHT_GRAY, bc=GOLD)

# ══════════════════════════════════════════════════════════════
# SLIDE 11 — Limitations & Next Steps
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_chrome(s, "10", "Next Steps & Path to Production")

# Left: current limitations
rect(s, Inches(0.45), Inches(1.25), Inches(4.5), Inches(5.95), DARK_CARD)
rect(s, Inches(0.45), Inches(1.25), Inches(4.5), Inches(0.06), RED_LIGHT)
txt(s, "CURRENT LIMITATIONS", Inches(0.6), Inches(1.38),
    Inches(4.2), Inches(0.3), size=11, bold=True, color=RED_LIGHT)
limits = [
    ("Inference speed", "~2s per image — too slow for real-time try-on"),
    ("Body types",      "Model trained on specific body proportions; may not generalise to all sizes"),
    ("Occlusion",       "Complex backgrounds or accessories can confuse the segmentation"),
    ("Lower-body",      "Pipeline handles tops only — no pants, skirts, or full outfits yet"),
    ("Fine details",    "Small patterns (fine stripes, micro-prints) occasionally blurred"),
]
ly = Inches(1.82)
for cat, desc in limits:
    txt(s, f"⚠  {cat}:", Inches(0.6), ly, Inches(4.2), Inches(0.3),
        size=11, bold=True, color=RED_LIGHT)
    txt(s, desc, Inches(0.6), ly + Inches(0.3), Inches(4.2), Inches(0.42),
        size=11, color=LIGHT_GRAY)
    ly += Inches(0.85)

# Right: roadmap
rect(s, Inches(5.25), Inches(1.25), Inches(7.7), Inches(5.95), DARK_CARD)
rect(s, Inches(5.25), Inches(1.25), Inches(7.7), Inches(0.06), GREEN)
txt(s, "ROADMAP TO MVP", Inches(5.4), Inches(1.38),
    Inches(7.4), Inches(0.3), size=11, bold=True, color=GREEN)

phases = [
    ("Phase 1  ·  Immediate (1–2 months)", ACCENT, [
        "Wrap inference pipeline in a REST API (FastAPI / Flask)",
        "Add batch processing — handle 100+ try-on requests in parallel",
        "Build a simple web UI for stakeholder demos",
    ]),
    ("Phase 2  ·  Short-term (3–4 months)", GOLD, [
        "Fine-tune HR-VITON on brand-specific garment catalog",
        "Optimise model with TensorRT / ONNX for <500ms inference",
        "Extend to lower-body garments (pants, skirts)",
    ]),
    ("Phase 3  ·  Production (6–12 months)", GREEN, [
        "Deploy on scalable cloud (AWS/GCP) with auto-scaling",
        "Integrate with e-commerce platform (Shopify / Magento plugin)",
        "Add user size/body-shape personalisation",
        "Target: <200ms latency · 99.9% uptime · 1M+ try-ons/day",
    ]),
]
py = Inches(1.85)
for phase_title, pc, items in phases:
    rect(s, Inches(5.35), py, Inches(0.06), Inches(0.35), pc)
    txt(s, phase_title, Inches(5.5), py, Inches(7.3), Inches(0.35),
        size=11, bold=True, color=pc)
    py += Inches(0.38)
    bullets(s, items, Inches(5.5), py, Inches(7.2), Inches(len(items) * 0.38 + 0.1),
            size=11, color=LIGHT_GRAY, bc=pc)
    py += len(items) * Inches(0.38) + Inches(0.25)

# ══════════════════════════════════════════════════════════════
# SLIDE 12 — Appendix A: Technical Deep-Dive
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_chrome(s, "A1", "Appendix — Technical Architecture Details")

txt(s, "For the technically curious — how HR-VITON actually works under the hood.",
    Inches(0.5), Inches(1.2), Inches(12.5), Inches(0.4),
    size=12, color=MID_GRAY, italic=True)

# Stage 1 box
rect(s, Inches(0.45), Inches(1.7), Inches(5.85), Inches(5.5), DARK_CARD)
rect(s, Inches(0.45), Inches(1.7), Inches(5.85), Inches(0.07), ACCENT)
txt(s, "STAGE 1 — ConditionGenerator (tocg)", Inches(0.6), Inches(1.83),
    Inches(5.6), Inches(0.35), size=12, bold=True, color=ACCENT)
stage1 = [
    "Inputs: cloth image (3ch) + cloth mask (1ch) → 4-channel tensor",
    "Inputs: parse-agnostic map (13ch) + DensePose UV (3ch) → 16-channel tensor",
    "Outputs a 5-scale coarse-to-fine flow field (pixel-level warp map)",
    "Also predicts a 13-channel segmentation map (where each body part is)",
    "Flow is applied at full 768×1024 resolution via F.grid_sample",
    "remove_overlap() prevents garment from bleeding over face/arms",
    "Architecture: dual-encoder UNet with cross-attention Feature Fusion Block",
]
bullets(s, stage1, Inches(0.6), Inches(2.3),
        Inches(5.5), Inches(4.5), size=11, color=LIGHT_GRAY)

# Stage 2 box
rect(s, Inches(6.6), Inches(1.7), Inches(6.4), Inches(5.5), DARK_CARD)
rect(s, Inches(6.6), Inches(1.7), Inches(6.4), Inches(0.07), GOLD)
txt(s, "STAGE 2 — SPADEGenerator", Inches(6.75), Inches(1.83),
    Inches(6.15), Inches(0.35), size=12, bold=True, color=GOLD)
stage2 = [
    "Inputs: agnostic (3ch) + DensePose (3ch) + warped cloth (3ch) = 9 channels",
    "Conditioned on a 7-channel body parse map throughout the decoder",
    "ALIAS blocks: spatially adaptive normalisation at each resolution level",
    "ALIAS handles misalignment — cloth/body don't need to be pixel-perfect",
    "Decoder upsamples 4× to produce final 768×1024 RGB image",
    "Architecture: SPADEGenerator with spectral normalisation (alias instance norm)",
    "Parameter count: ~90M (ConditionGenerator) + ~60M (SPADEGenerator)",
]
bullets(s, stage2, Inches(6.75), Inches(2.3),
        Inches(6.1), Inches(4.5), size=11, color=LIGHT_GRAY, bc=GOLD)

# ══════════════════════════════════════════════════════════════
# SLIDE 13 — Appendix B: Full Metrics & Training Logs
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank)
slide_chrome(s, "A2", "Appendix — Extended Results & Training Analysis")

# Left — training progression chart (text-based)
rect(s, Inches(0.45), Inches(1.25), Inches(5.8), Inches(5.95), DARK_CARD)
txt(s, "v3 TRAINING LOSS — GMM STAGE (20 epochs)", Inches(0.6), Inches(1.4),
    Inches(5.5), Inches(0.3), size=11, bold=True, color=ACCENT)

loss_data = [
    ("Ep 1",  "0.187", "0.148"), ("Ep 4",  "0.121", "0.118"),
    ("Ep 7",  "0.114", "0.113"), ("Ep 10", "0.103", "0.101"),
    ("Ep 13", "0.095", "0.095"), ("Ep 16", "0.092", "0.093"),
]
col_heads = ["Epoch", "Train Loss", "Val Loss"]
hx = [Inches(0.6), Inches(2.0), Inches(3.5)]
ty_h = Inches(1.82)
rect(s, Inches(0.55), ty_h, Inches(5.6), Inches(0.3), ACCENT)
for hdr, hxx in zip(col_heads, hx):
    txt(s, hdr, hxx, ty_h + Inches(0.02), Inches(1.3), Inches(0.26),
        size=9.5, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
for i, (ep, tr, vl) in enumerate(loss_data):
    ry = ty_h + Inches(0.3) + i * Inches(0.35)
    bg = DARK_CARD if i % 2 == 0 else RGBColor(0x0E,0x1E,0x42)
    rect(s, Inches(0.55), ry, Inches(5.6), Inches(0.33), bg)
    for val, hxx in zip([ep, tr, vl], hx):
        txt(s, val, hxx, ry + Inches(0.03), Inches(1.3), Inches(0.28),
            size=10, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

txt(s, "Observation: train and val loss track closely → no overfitting. "
       "Gap < 0.003 throughout — regularisation (dropout + weight decay) working.",
    Inches(0.6), Inches(4.25), Inches(5.5), Inches(0.75),
    size=11, color=GOLD, italic=True)

txt(s, "WHY SSIM STAYED LOW DESPITE GOOD TRAINING LOSS", Inches(0.6), Inches(5.15),
    Inches(5.5), Inches(0.3), size=10, bold=True, color=RED_LIGHT)
txt(s, "Training loss (mask alignment) improved steadily, but SSIM measures "
       "full image structural similarity — the model learned to place the mask "
       "correctly but couldn't synthesise realistic texture without 10x more data.",
    Inches(0.6), Inches(5.5), Inches(5.5), Inches(1.4),
    size=11, color=LIGHT_GRAY)

# Right — hyperparameter summary
rect(s, Inches(6.55), Inches(1.25), Inches(6.4), Inches(5.95), DARK_CARD)
txt(s, "HYPERPARAMETER SUMMARY", Inches(6.7), Inches(1.4),
    Inches(6.15), Inches(0.3), size=11, bold=True, color=GOLD)

params = [
    ("Optimizer",         "Adam (β₁=0.9, β₂=0.999)"),
    ("Learning rate",     "1×10⁻⁴  (v2) · OneCycleLR max 5×10⁻⁴ (v3)"),
    ("Batch size",        "4 (v2)  ·  8 (v3)"),
    ("Training pairs",    "600 (v2)  ·  5,000 (v3)"),
    ("Max epochs",        "20/20 per stage (v2)  ·  30/30 (v3)"),
    ("Early stopping",    "Patience = 5 (v2)  ·  3 (v3)"),
    ("Grad clipping",     "max_norm = 1.0"),
    ("Weight decay",      "1×10⁻⁴"),
    ("GMM loss (v2)",     "L1(mask) + 0.01·‖θ‖₂"),
    ("GMM loss (v3)",     "L1(mask) + 0.5·L1(appearance) + 0.001·TV"),
    ("TOM loss",          "L1 + 0.25·VGG + 0.5·(1−SSIM) + 0.01·TV"),
    ("Dropout",           "0.1 / 0.2 / 0.3 at encoder levels 1/2/3"),
    ("Augmentation",      "Horizontal flip (p=0.5) + colour jitter"),
    ("Precision",         "FP32 (v2 CPU)  ·  AMP FP16 (v3 GPU)"),
]
py = Inches(1.85)
for cat, val in params:
    txt(s, cat + ":", Inches(6.7), py, Inches(1.9), Inches(0.3),
        size=10, bold=True, color=ACCENT)
    txt(s, val, Inches(8.65), py, Inches(4.1), Inches(0.3),
        size=10, color=LIGHT_GRAY)
    py += Inches(0.38)

# ── Save ────────────────────────────────────────────────────────
out = "VirtualTryOn_Presentation.pptx"
prs.save(out)
print(f"Saved → {out}  ({len(prs.slides)} slides)")
